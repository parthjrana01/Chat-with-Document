import streamlit as st
from PyPDF2 import PdfReader
import docx
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import pypandoc
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import google.generativeai as genai

with open('style.css') as f:
    st.markdown(f'<style>{f.read()}</style>', unsafe_allow_html=True)

# Rest of your Streamlit code
# st.title("My Styled Streamlit App")
# st.header("This is a header with custom font")

# Load environment variables
load_dotenv()

# Configure Google Generative AI
genai.configure(api_key=os.getenv("AIzaSyAbGYl1RWYku3ntot7fWhRXjwZlwNOJzvc"))

# Function to get text from a DOCX file
def get_doc_text(doc_files):
    doc = Document(doc_files)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

# Function to get text from a PPTX file
def get_ppt_text(ppt_files):
    try:
        prs = Presentation(ppt_files)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error("Error: " + str(e))

# Function to get text from HTML files
def get_html_text(html_files):
    text = ""
    for html_file in html_files:
        with open(html_file, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, 'html.parser')
            text += soup.get_text()
    return text

# Function to get text from LaTeX files
def get_latex_text(latex_files):
    text = ""
    for latex_file in latex_files:
        output = pypandoc.convert_file(latex_file, 'plain')
        text += output
    return text

# Function to get text from PDF files
def get_pdf_text(pdf_docs):
    text = ""
    pdf_reader = PdfReader(pdf_docs)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

# Function to parse text from a file
def parse_text(file):
    text = file.getvalue().decode("utf-8")
    return text

# Function to get text from uploaded documents
def get_text_from_documents(pdf_docs):
    content = ""
    for file in pdf_docs:
        if file is not None:
            if file.name.endswith('.pdf'):
                content += get_pdf_text(file)
            elif file.name.endswith('.docx') or file.name.endswith('.doc'):
                content += get_doc_text(file)
            elif file.name.endswith('.pptx') or file.name.endswith('.ppt'):
                content += get_ppt_text(file)
            elif file.name.endswith('.html'):
                content += get_html_text(file)
            elif file.name.endswith('.tex'):
                content += get_latex_text(file)
            elif file.name.endswith('.txt'):
                content += parse_text(file)
    return content

# Function to split text into chunks
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

# Function to create a vector store
def get_vector_store(text_chunks):
    if not text_chunks:
        st.warning("No text chunks found.")
        return
    
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

# Function to create a conversational chain
def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.1)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

# Function to handle user input
def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()
    response = chain(
        {"input_documents": docs, "question": user_question},
        return_only_outputs=True
    )
    output_text = response.get("output_text", "No answer available")
    # Replace bullet points with line breaks
    output_text = output_text.replace('•', '\n•')
    # Add a border around the entire response
    st.markdown(f"<div style='border: 1px solid #ccc; padding: 10px;'>🤖: {output_text}</div>", unsafe_allow_html=True)
    return output_text

# Streamlit app
def main():
    # st.set_page_config("Chat PDF")
    st.title("Chat with PDF using Gemini💁")
    st.markdown("Ask questions about the uploaded documents.")

    # Load chat history from session state
    chat_history = st.session_state.get("chat_history", [])

    # User input
    user_question = st.text_input("Ask a Question from the PDF Files", help="Enter your question here.")

    # Sidebar
    with st.sidebar:
        st.title("")
        st.write("Upload your PDF files and click on the 'Submit & Process' button.")
        pdf_docs = st.file_uploader("Select PDF Files", accept_multiple_files=True)
        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                # Clear existing session state
                st.session_state.chat_history = []
                # Process documents
                raw_text = get_text_from_documents(pdf_docs)
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks)
                st.success("Processing done!")

    # Process user input and update chat history
    if user_question:
        # Process user input and get response
        response = user_input(user_question)

        # Update chat history
        chat_history.append("<br>")
        chat_history.append(f"💁: {user_question}")
        chat_history.append(f"🤖: {response}")

        # Save updated chat history to session state
        st.session_state.chat_history = chat_history

    # Display chat history as paragraphs with scrolling
    st.markdown("<div style='height:400px;overflow-y:scroll;'>" + "<br>".join(chat_history) + "</div>", unsafe_allow_html=True)

if __name__ == "__main__":
    main()